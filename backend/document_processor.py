"""Document processing for different file types"""
import os
import re
import logging
import shutil
import hashlib
import json
import zipfile
import base64
import unicodedata
from typing import List, Dict, Optional, Tuple
from io import BytesIO
import httpx
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process different document types and extract text"""
    
    CHUNK_SIZE = 1000
    CHUNK_OVERLAP = 200

    @classmethod
    def _chunk_settings(
        cls,
        *,
        file_size_bytes: int = 0,
        page_count: int = 0,
        used_ocr: bool = False,
        text_length: int = 0,
    ) -> Tuple[int, int]:
        """Use larger chunks for very large documents to reduce embedding load."""
        size_mb = max(file_size_bytes / (1024 * 1024), 0.0)

        chunk_size = cls.CHUNK_SIZE
        overlap = cls.CHUNK_OVERLAP

        if text_length >= 500_000 or page_count >= 60 or size_mb >= 15:
            chunk_size = 1600
            overlap = 240
        if used_ocr or text_length >= 1_000_000 or page_count >= 120 or size_mb >= 30:
            chunk_size = 2200
            overlap = 260
        if text_length >= 2_000_000 or page_count >= 180 or size_mb >= 45:
            chunk_size = 2800
            overlap = 320
        if used_ocr and (text_length >= 3_000_000 or page_count >= 240 or size_mb >= 60):
            chunk_size = 3600
            overlap = 360

        return chunk_size, overlap

    @staticmethod
    def _ocr_render_scale(page_count: int, file_size_bytes: int) -> float:
        """Reduce OCR image resolution for very large scanned PDFs to avoid timeouts."""
        size_mb = max(file_size_bytes / (1024 * 1024), 0.0)
        if page_count >= 150 or size_mb >= 60:
            return 1.0
        if page_count >= 100 or size_mb >= 40:
            return 1.15
        if page_count >= 60 or size_mb >= 25:
            return 1.3
        if page_count >= 25 or size_mb >= 10:
            return 1.55
        return 2.0

    @staticmethod
    def _resolve_tesseract_command() -> Optional[str]:
        """Resolve the Tesseract executable for OCR."""
        configured = os.environ.get("TESSERACT_CMD")
        if configured and os.path.exists(configured):
            return configured

        return shutil.which("tesseract")
    
    @staticmethod
    def clean_text(text: str) -> str:
        """Clean and normalize extracted text"""
        text = unicodedata.normalize("NFKC", text or "")
        text = text.replace("\u00a0", " ")
        # Remove multiple newlines and spaces
        text = re.sub(r'\n+', '\n', text)
        text = re.sub(r' +', ' ', text)
        text = text.strip()
        return text

    @classmethod
    def clean_ocr_text(cls, text: str) -> str:
        """Normalize OCR-heavy text while keeping useful structure."""
        if not text:
            return ""

        text = text.replace("\x0c", "\n")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"(?<=\w)-\s*\n\s*(?=\w)", "", text)
        text = re.sub(r"(?<=\w)\s*\n\s*(?=\w)", " ", text)
        text = re.sub(r"[|]{2,}", "|", text)
        text = re.sub(r"[_]{2,}", "_", text)

        cleaned_lines = []
        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                continue

            alnum_count = sum(char.isalnum() for char in line)
            if alnum_count == 0:
                continue

            symbol_count = sum(not char.isalnum() and not char.isspace() for char in line)
            if len(line) >= 8 and (symbol_count / max(len(line), 1)) > 0.45:
                continue

            cleaned_lines.append(line)

        return cls.clean_text("\n".join(cleaned_lines))

    @staticmethod
    def _text_quality_metrics(text: str) -> Dict[str, float]:
        """Estimate extraction quality for OCR or weak PDF text."""
        normalized = (text or "").strip()
        if not normalized:
            return {
                "quality_score": 0.0,
                "word_count": 0,
                "alpha_ratio": 0.0,
                "digit_ratio": 0.0,
                "noise_ratio": 1.0,
                "avg_word_length": 0.0,
            }

        letters = sum(char.isalpha() for char in normalized)
        digits = sum(char.isdigit() for char in normalized)
        spaces = sum(char.isspace() for char in normalized)
        total_chars = len(normalized)
        noise_chars = max(total_chars - letters - digits - spaces, 0)

        words = re.findall(r"\b[\w/-]+\b", normalized)
        word_count = len(words)
        avg_word_length = (
            sum(len(word) for word in words) / word_count
            if word_count
            else 0.0
        )

        alpha_ratio = letters / total_chars
        digit_ratio = digits / total_chars
        noise_ratio = noise_chars / total_chars

        score = 0.0
        if word_count >= 20:
            score += 0.35
        elif word_count >= 8:
            score += 0.2
        else:
            score += 0.08

        score += min(alpha_ratio, 0.55)
        score += max(0.0, 0.18 - noise_ratio)
        if 2.5 <= avg_word_length <= 12:
            score += 0.08
        if digit_ratio > 0.35:
            score -= 0.08

        return {
            "quality_score": round(max(0.0, min(score, 1.0)), 3),
            "word_count": word_count,
            "alpha_ratio": round(alpha_ratio, 3),
            "digit_ratio": round(digit_ratio, 3),
            "noise_ratio": round(noise_ratio, 3),
            "avg_word_length": round(avg_word_length, 3),
        }

    @classmethod
    def _should_force_ocr(cls, extracted_text: str, page_count: int) -> bool:
        """Decide when a PDF text layer is too weak and OCR should be attempted."""
        metrics = cls._text_quality_metrics(extracted_text)
        if metrics["word_count"] == 0:
            return True
        if page_count >= 2 and metrics["word_count"] < page_count * 25:
            return True
        return metrics["quality_score"] < 0.32

    @staticmethod
    def _merge_text_versions(primary_text: str, secondary_text: str) -> Tuple[str, Dict[str, float]]:
        """Choose the cleaner text version using quality metrics."""
        primary_metrics = DocumentProcessor._text_quality_metrics(primary_text)
        secondary_metrics = DocumentProcessor._text_quality_metrics(secondary_text)

        if secondary_metrics["quality_score"] > primary_metrics["quality_score"] + 0.08:
            return secondary_text, secondary_metrics
        return primary_text, primary_metrics

    @staticmethod
    def _looks_like_content_image_url(url: str) -> bool:
        """Keep content images and skip obvious logos, icons, and tiny assets."""
        normalized = (url or "").strip()
        if not normalized:
            return False

        parsed = urlparse(normalized)
        if parsed.scheme not in {"http", "https"}:
            return False

        lowered = normalized.lower()
        if any(token in lowered for token in ("logo", "icon", "sprite", "favicon", "placeholder", "avatar")):
            return False
        if lowered.endswith((".svg", ".ico", ".gif")):
            return False
        return True

    @staticmethod
    def _prepare_image_asset(image_bytes: bytes, filename: str, alt: str = "") -> Optional[Dict]:
        """Normalize extracted document images into lightweight chat-ready assets."""
        if not image_bytes:
            return None

        try:
            from PIL import Image
        except ImportError:
            return None

        try:
            with Image.open(BytesIO(image_bytes)) as image:
                width, height = image.size
                if width < 180 or height < 180:
                    return None

                working = image.copy()
                working.thumbnail((1280, 1280))
                output = BytesIO()

                if working.mode in ("RGBA", "LA", "P"):
                    working.save(output, format="PNG", optimize=True)
                    content_type = "image/png"
                    extension = "png"
                else:
                    if working.mode != "RGB":
                        working = working.convert("RGB")
                    working.save(output, format="JPEG", quality=84, optimize=True)
                    content_type = "image/jpeg"
                    extension = "jpg"

                return {
                    "filename": f"{os.path.splitext(filename)[0]}.{extension}",
                    "content": output.getvalue(),
                    "content_type": content_type,
                    "alt": (alt or "").strip(),
                }
        except Exception as exc:
            logger.debug("Skipping image asset %s due to processing error: %s", filename, exc)
            return None

    @classmethod
    def extract_html_images(
        cls,
        soup: BeautifulSoup,
        base_url: str,
        *,
        source_title: str = "",
        max_images: int = 4,
    ) -> List[Dict]:
        """Extract a few useful absolute image URLs from a webpage."""
        if not soup or not base_url or max_images <= 0:
            return []

        images = []
        seen = set()

        def add_image(url_value: str, alt: str = "") -> None:
            if len(images) >= max_images:
                return
            absolute_url = urljoin(base_url, (url_value or "").strip())
            normalized_key = absolute_url.split("#")[0]
            if not cls._looks_like_content_image_url(absolute_url) or normalized_key in seen:
                return
            seen.add(normalized_key)
            images.append(
                {
                    "url": absolute_url,
                    "alt": cls.clean_text(alt or source_title)[:160],
                    "source_title": source_title,
                    "source_url": base_url,
                    "origin": "website",
                }
            )

        for selector in (
            ('meta[property="og:image"]', "content"),
            ('meta[name="twitter:image"]', "content"),
        ):
            for tag in soup.select(selector[0]):
                add_image(tag.get(selector[1], ""))
                if len(images) >= max_images:
                    return images

        main_content = soup.find(["main", "article"]) or soup.find("body")
        for image in (main_content.find_all("img") if main_content else []):
            src = image.get("src") or image.get("data-src") or image.get("data-lazy-src") or ""
            width = int(image.get("width") or 0) if str(image.get("width") or "").isdigit() else 0
            height = int(image.get("height") or 0) if str(image.get("height") or "").isdigit() else 0
            if width and width < 180:
                continue
            if height and height < 180:
                continue
            add_image(src, image.get("alt", ""))
            if len(images) >= max_images:
                break

        return images

    @classmethod
    def _build_ocr_engines(cls):
        """Initialize OCR engines once so page/image OCR can share the same path."""
        tesseract_cmd = cls._resolve_tesseract_command()
        pytesseract = None
        rapid_ocr = None
        ocr_engine = None

        if tesseract_cmd:
            try:
                import pytesseract as pytesseract_module

                pytesseract_module.pytesseract.tesseract_cmd = tesseract_cmd
                pytesseract = pytesseract_module
                ocr_engine = "tesseract"
            except Exception as exc:
                logger.warning("Tesseract OCR unavailable, falling back to RapidOCR: %s", exc)

        if not pytesseract:
            try:
                from rapidocr_onnxruntime import RapidOCR

                rapid_ocr = RapidOCR()
                ocr_engine = "rapidocr"
            except Exception as exc:
                logger.warning("RapidOCR unavailable: %s", exc)

        return pytesseract, rapid_ocr, ocr_engine

    @classmethod
    def _ocr_pil_image(cls, image, pytesseract=None, rapid_ocr=None) -> str:
        """OCR a PIL image using whichever OCR engine is available."""
        if image is None:
            return ""

        page_text = ""
        if pytesseract:
            page_text = pytesseract.image_to_string(image, config="--psm 6")
        elif rapid_ocr:
            import numpy as np

            ocr_result, _ = rapid_ocr(np.array(image))
            page_text = "\n".join(item[1] for item in ocr_result) if ocr_result else ""

        return cls.clean_ocr_text(page_text)

    @classmethod
    def _extract_pdf_page_layout_text(cls, page) -> str:
        """Preserve page text in reading order using PyMuPDF's layout-aware dict output."""
        text_dict = page.get_text("dict")
        block_texts = []

        for block in sorted(text_dict.get("blocks", []), key=lambda item: (item.get("bbox", [0, 0])[1], item.get("bbox", [0, 0])[0])):
            if block.get("type") != 0:
                continue

            line_texts = []
            for line in block.get("lines", []):
                spans = [span.get("text", "") for span in line.get("spans", []) if span.get("text", "").strip()]
                joined = cls.clean_text("".join(spans))
                if joined:
                    line_texts.append(joined)

            if line_texts:
                block_texts.append("\n".join(line_texts))

        return cls.clean_text("\n\n".join(block_texts))

    @classmethod
    def _extract_pdf_page_table_lines(cls, page, page_number: int) -> List[str]:
        """Extract table rows from one PDF page as retrieval-friendly facts."""
        if not hasattr(page, "find_tables"):
            return []

        lines = []
        try:
            table_finder = page.find_tables()
        except Exception as exc:
            logger.debug("PyMuPDF table detection failed on page %s: %s", page_number, exc)
            return []

        for table_index, table in enumerate(getattr(table_finder, "tables", []) or [], start=1):
            try:
                rows = table.extract()
            except Exception as exc:
                logger.debug("PyMuPDF table extraction failed on page %s: %s", page_number, exc)
                continue

            lines.extend(
                cls._format_table_matrix(rows, table_label=f"PDF table {table_index} on page {page_number}")
            )

        return lines

    @classmethod
    def _describe_image_with_gemini(cls, image_bytes: bytes, page_number: int) -> str:
        """Describe a PDF image with Gemini when configured so image semantics enter the index."""
        api_key = os.environ.get("GEMINI_API_KEY", "").strip()
        if not api_key or not image_bytes:
            return ""

        model = os.environ.get("GEMINI_MODEL", "gemini-1.5-flash").strip() or "gemini-1.5-flash"
        mime_type = "image/png"
        if image_bytes.startswith(b"\xff\xd8"):
            mime_type = "image/jpeg"
        elif image_bytes.startswith(b"RIFF") and b"WEBP" in image_bytes[:16]:
            mime_type = "image/webp"

        payload = {
            "contents": [
                {
                    "parts": [
                        {
                            "text": (
                                "Describe this document image for retrieval. Focus only on visible factual content "
                                "such as captions, charts, posters, labels, tables, campus names, programme names, "
                                "or other academic details. Keep it under 80 words and do not invent missing details."
                            )
                        },
                        {
                            "inline_data": {
                                "mime_type": mime_type,
                                "data": base64.b64encode(image_bytes).decode("ascii"),
                            }
                        },
                    ]
                }
            ],
            "generationConfig": {"temperature": 0.1, "topP": 0.8, "maxOutputTokens": 160},
        }

        try:
            response = httpx.post(
                f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent",
                params={"key": api_key},
                json=payload,
                timeout=40.0,
            )
            response.raise_for_status()
            data = response.json()
            parts = data.get("candidates", [{}])[0].get("content", {}).get("parts", [])
            text = " ".join(part.get("text", "").strip() for part in parts if part.get("text"))
            text = cls.clean_text(text)
            if not text:
                return ""
            return f"Image description on page {page_number}: {text}"
        except Exception as exc:
            logger.debug("Gemini image description skipped on page %s: %s", page_number, exc)
            return ""

    @classmethod
    def _extract_pdf_page_image_details(
        cls,
        pdf,
        page,
        page_number: int,
        *,
        pytesseract=None,
        rapid_ocr=None,
    ) -> Tuple[List[Dict], List[str]]:
        """Extract useful PDF images and convert their visible meaning into searchable text."""
        details = []
        extracted_images = []

        try:
            from PIL import Image, ImageFilter, ImageOps
        except ImportError:
            Image = ImageFilter = ImageOps = None

        seen_hashes = set()
        for image_number, image_info in enumerate(page.get_images(full=True), start=1):
            try:
                image_data = pdf.extract_image(image_info[0])
                raw_bytes = image_data.get("image")
                if not raw_bytes:
                    continue

                digest = hashlib.sha1(raw_bytes).hexdigest()
                if digest in seen_hashes:
                    continue
                seen_hashes.add(digest)

                prepared = cls._prepare_image_asset(
                    raw_bytes,
                    filename=f"page-{page_number}-image-{image_number}.{image_data.get('ext', 'png')}",
                    alt=f"Document image from page {page_number}",
                )
                if prepared:
                    extracted_images.append(prepared)

                image_notes = []
                if Image and ImageOps:
                    try:
                        pil_image = Image.open(BytesIO(raw_bytes))
                        pil_image = ImageOps.grayscale(pil_image)
                        pil_image = ImageOps.autocontrast(pil_image)
                        if ImageFilter:
                            pil_image = pil_image.filter(ImageFilter.SHARPEN)
                        image_ocr_text = cls._ocr_pil_image(pil_image, pytesseract=pytesseract, rapid_ocr=rapid_ocr)
                        if image_ocr_text and len(image_ocr_text) >= 12:
                            image_notes.append(f"Image text on page {page_number}: {image_ocr_text}")
                    except Exception as exc:
                        logger.debug("Image OCR skipped on page %s image %s: %s", page_number, image_number, exc)

                gemini_description = cls._describe_image_with_gemini(raw_bytes, page_number)
                if gemini_description:
                    image_notes.append(gemini_description)

                if image_notes:
                    details.extend(image_notes)
                    if prepared:
                        prepared["alt"] = image_notes[0][:160]
            except Exception as exc:
                logger.debug("Skipping PDF image detail extraction on page %s: %s", page_number, exc)

        return extracted_images, details

    @classmethod
    def _extract_pdf_images(cls, file_content: bytes, max_images: int = 4) -> List[Dict]:
        """Extract a few meaningful images from a PDF."""
        try:
            import fitz
        except ImportError:
            return []

        extracted = []
        seen_hashes = set()

        try:
            pdf = fitz.open(stream=file_content, filetype="pdf")
        except Exception as exc:
            logger.debug("PDF image extraction skipped: %s", exc)
            return []

        try:
            page_limit = min(pdf.page_count, 24)
            for page_index in range(page_limit):
                page = pdf.load_page(page_index)
                for image_number, image_info in enumerate(page.get_images(full=True), start=1):
                    if len(extracted) >= max_images:
                        return extracted

                    try:
                        image_data = pdf.extract_image(image_info[0])
                        raw_bytes = image_data.get("image")
                        if not raw_bytes:
                            continue
                        digest = hashlib.sha1(raw_bytes).hexdigest()
                        if digest in seen_hashes:
                            continue

                        prepared = cls._prepare_image_asset(
                            raw_bytes,
                            filename=f"page-{page_index + 1}-image-{image_number}.{image_data.get('ext', 'png')}",
                            alt=f"Document image from page {page_index + 1}",
                        )
                        if not prepared:
                            continue

                        extracted.append(prepared)
                        seen_hashes.add(digest)
                    except Exception as exc:
                        logger.debug("Skipping PDF image on page %s: %s", page_index + 1, exc)
                        continue
        finally:
            pdf.close()

        return extracted

    @staticmethod
    def _format_table_matrix(table_rows: List[List], table_label: str = "Table") -> List[str]:
        """Convert a rectangular table into retrieval-friendly row statements."""
        normalized_rows = []
        for row in table_rows or []:
            normalized_row = []
            for cell in row or []:
                value = DocumentProcessor.clean_text(str(cell or ""))
                normalized_row.append(value)
            if any(normalized_row):
                normalized_rows.append(normalized_row)

        if not normalized_rows:
            return []

        lines = [table_label]
        header = normalized_rows[0]

        def looks_like_header(row: List[str]) -> bool:
            filled = [cell for cell in row if cell]
            if len(filled) < 2:
                return False
            joined = " ".join(filled).lower()
            header_markers = {
                "parameter",
                "details",
                "particulars",
                "description",
                "course",
                "code",
                "credits",
                "semester",
                "subject",
                "component",
            }
            return any(marker in joined for marker in header_markers)

        has_header = looks_like_header(header)
        data_rows = normalized_rows[1:] if has_header else normalized_rows

        for row_index, row in enumerate(data_rows, start=1):
            if has_header:
                pairs = []
                for column_index, cell in enumerate(row):
                    if not cell:
                        continue
                    heading = header[column_index] if column_index < len(header) and header[column_index] else f"Column {column_index + 1}"
                    pairs.append(f"{heading}: {cell}")
                if pairs:
                    lines.append(f"Table row {row_index}: " + "; ".join(pairs))
            elif len([cell for cell in row if cell]) == 2:
                left, right = [cell for cell in row if cell]
                lines.append(f"Table row {row_index}: {left}: {right}")
            else:
                joined = " | ".join(cell for cell in row if cell)
                if joined:
                    lines.append(f"Table row {row_index}: {joined}")

        return lines if len(lines) > 1 else []

    @classmethod
    def _extract_pdf_tables_text(cls, file_content: bytes) -> str:
        """Extract ruled PDF tables with PyMuPDF and preserve cell relationships."""
        try:
            import fitz
        except ImportError:
            return ""

        table_lines = []
        try:
            pdf = fitz.open(stream=file_content, filetype="pdf")
        except Exception as exc:
            logger.debug("PDF table extraction skipped: %s", exc)
            return ""

        try:
            for page_index, page in enumerate(pdf, start=1):
                if not hasattr(page, "find_tables"):
                    return ""
                try:
                    table_finder = page.find_tables()
                except Exception as exc:
                    logger.debug("PyMuPDF table detection failed on page %s: %s", page_index, exc)
                    continue

                for table_index, table in enumerate(getattr(table_finder, "tables", []) or [], start=1):
                    try:
                        rows = table.extract()
                    except Exception as exc:
                        logger.debug("PyMuPDF table extraction failed on page %s: %s", page_index, exc)
                        continue

                    table_lines.extend(
                        cls._format_table_matrix(rows, table_label=f"PDF table {table_index} on page {page_index}")
                    )
        finally:
            pdf.close()

        return cls.clean_text("\n".join(table_lines))

    @classmethod
    def _extract_docx_images(cls, file_content: bytes, max_images: int = 4) -> List[Dict]:
        """Extract embedded images from a DOCX document."""
        try:
            from docx import Document
        except ImportError:
            return []

        extracted = []
        seen_hashes = set()

        try:
            doc = Document(BytesIO(file_content))
            for index, rel in enumerate(doc.part.rels.values(), start=1):
                if len(extracted) >= max_images or "image" not in rel.reltype:
                    continue

                raw_bytes = getattr(rel.target_part, "blob", None)
                if not raw_bytes:
                    continue

                digest = hashlib.sha1(raw_bytes).hexdigest()
                if digest in seen_hashes:
                    continue

                original_name = os.path.basename(str(getattr(rel.target_part, "partname", f"docx-image-{index}.png")))
                prepared = cls._prepare_image_asset(
                    raw_bytes,
                    filename=original_name or f"docx-image-{index}.png",
                    alt=f"Embedded document image {index}",
                )
                if not prepared:
                    continue

                extracted.append(prepared)
                seen_hashes.add(digest)
        except Exception as exc:
            logger.debug("DOCX image extraction skipped: %s", exc)

        return extracted
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if not text:
            return []
        
        chunks = []
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        current_chunk = ""
        for sentence in sentences:
            if len(current_chunk) + len(sentence) <= chunk_size:
                current_chunk += " " + sentence if current_chunk else sentence
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    # Keep some overlap
                    words = current_chunk.split()
                    overlap_words = words[-overlap//5:] if len(words) > overlap//5 else words
                    current_chunk = " ".join(overlap_words) + " " + sentence
                else:
                    # Sentence is too long, split by words
                    words = sentence.split()
                    for i in range(0, len(words), chunk_size//5):
                        chunk_words = words[i:i + chunk_size//5]
                        chunks.append(" ".join(chunk_words))
                    current_chunk = ""
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    @classmethod
    def process_pdf(cls, file_content: bytes) -> Dict:
        """Extract text from PDF file"""
        try:
            import fitz
            from PIL import Image, ImageFilter, ImageOps

            pdf = fitz.open(stream=file_content, filetype="pdf")
            try:
                page_count = pdf.page_count
                pytesseract, rapid_ocr, detected_ocr_engine = cls._build_ocr_engines()
                page_sections = []
                page_texts = []
                all_images = []
                used_ocr = False
                ocr_engine = None

                render_scale = cls._ocr_render_scale(page_count, len(file_content))

                for page_index in range(page_count):
                    page_number = page_index + 1
                    page = pdf.load_page(page_index)

                    layout_text = cls._extract_pdf_page_layout_text(page)
                    table_lines = cls._extract_pdf_page_table_lines(page, page_number)
                    extracted_page_text = cls.clean_text(
                        "\n".join(
                            part for part in [layout_text, "\n".join(table_lines)] if part
                        )
                    )

                    extracted_metrics = cls._text_quality_metrics(extracted_page_text)
                    needs_ocr = (
                        not extracted_page_text
                        or extracted_metrics["word_count"] < 14
                        or extracted_metrics["quality_score"] < 0.24
                    )

                    ocr_page_text = ""
                    if needs_ocr and (pytesseract or rapid_ocr):
                        try:
                            pix = page.get_pixmap(matrix=fitz.Matrix(render_scale, render_scale), alpha=False)
                            image = Image.open(BytesIO(pix.tobytes("png")))
                            image = ImageOps.grayscale(image)
                            image = ImageOps.autocontrast(image)
                            image = image.filter(ImageFilter.SHARPEN)
                            ocr_page_text = cls._ocr_pil_image(
                                image,
                                pytesseract=pytesseract,
                                rapid_ocr=rapid_ocr,
                            )
                            if ocr_page_text:
                                used_ocr = True
                                ocr_engine = detected_ocr_engine
                        except Exception as exc:
                            logger.debug("Page OCR skipped on page %s: %s", page_number, exc)

                    merged_page_text, _ = cls._merge_text_versions(extracted_page_text, ocr_page_text)
                    page_texts.append(merged_page_text)

                    page_images, image_details = cls._extract_pdf_page_image_details(
                        pdf,
                        page,
                        page_number,
                        pytesseract=pytesseract,
                        rapid_ocr=rapid_ocr,
                    )
                    if page_images:
                        all_images.extend(page_images)

                    page_parts = [f"Page {page_number}"]
                    if merged_page_text:
                        page_parts.append(merged_page_text)
                    if image_details:
                        page_parts.append("\n".join(image_details))

                    page_section_text = cls.clean_text("\n\n".join(part for part in page_parts if part))
                    if page_section_text:
                        page_sections.append(page_section_text)

                extracted_text = cls.clean_text("\n\n".join(page_sections))

                if not extracted_text.strip() or cls._should_force_ocr(extracted_text, page_count):
                    ocr_result = cls.ocr_pdf(file_content, extracted_pages=page_texts)
                    if not ocr_result["success"]:
                        if not extracted_text.strip():
                            return ocr_result
                        full_text = extracted_text
                        quality_metrics = cls._text_quality_metrics(full_text)
                    else:
                        merged_text, quality_metrics = cls._merge_text_versions(
                            extracted_text,
                            cls.clean_ocr_text(ocr_result["text"]),
                        )
                        full_text = cls.clean_text(merged_text)
                        used_ocr = True
                        ocr_engine = ocr_result.get("ocr_engine") or ocr_engine
                else:
                    full_text = extracted_text
                    quality_metrics = cls._text_quality_metrics(full_text)

                chunk_size, overlap = cls._chunk_settings(
                    file_size_bytes=len(file_content),
                    page_count=page_count,
                    used_ocr=used_ocr,
                    text_length=len(full_text),
                )
                chunks = cls.chunk_text(full_text, chunk_size=chunk_size, overlap=overlap)
            finally:
                pdf.close()
            
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "images": all_images[:12],
                "page_count": page_count,
                "used_ocr": used_ocr,
                "ocr_engine": ocr_engine,
                "ocr_quality_score": quality_metrics.get("quality_score", 0.0),
                "extraction_quality": quality_metrics,
                "chunk_size": chunk_size,
            }
        except ImportError:
            logger.error("PDF support is unavailable because no PDF reader library is installed")
            return {
                "success": False,
                "error": "PDF support is not installed on the backend yet. Please install the PDF reader dependency.",
                "chunks": []
            }
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            return {"success": False, "error": str(e), "chunks": []}

    @classmethod
    def ocr_pdf(cls, file_content: bytes, extracted_pages: Optional[List[str]] = None) -> Dict:
        """Run OCR on scanned PDFs that do not contain a text layer."""
        tesseract_cmd = cls._resolve_tesseract_command()
        try:
            import fitz
            from PIL import Image, ImageFilter, ImageOps
        except ImportError as exc:
            logger.error("OCR dependencies missing: %s", exc)
            return {
                "success": False,
                "error": "OCR dependencies are not installed on the backend yet.",
                "chunks": []
            }

        try:
            pdf = fitz.open(stream=file_content, filetype="pdf")
            text_parts = []

            ocr_with_tesseract = False
            pytesseract = None
            rapid_ocr = None

            if tesseract_cmd:
                try:
                    import pytesseract as pytesseract_module
                    pytesseract_module.pytesseract.tesseract_cmd = tesseract_cmd
                    pytesseract = pytesseract_module
                    ocr_with_tesseract = True
                except Exception as exc:
                    logger.warning("Tesseract OCR unavailable, falling back to RapidOCR: %s", exc)

            if not pytesseract:
                try:
                    from rapidocr_onnxruntime import RapidOCR
                    rapid_ocr = RapidOCR()
                except Exception as exc:
                    logger.error("RapidOCR unavailable: %s", exc)
                    return {
                        "success": False,
                        "error": (
                            "This PDF appears to be scanned, but no OCR engine is available on the backend."
                        ),
                        "chunks": []
                    }

            render_scale = cls._ocr_render_scale(pdf.page_count, len(file_content))

            for page_number, page in enumerate(pdf, start=1):
                extracted_page_text = ""
                if extracted_pages and len(extracted_pages) >= page_number:
                    extracted_page_text = cls.clean_text(extracted_pages[page_number - 1] or "")

                extracted_metrics = cls._text_quality_metrics(extracted_page_text)
                needs_ocr = (
                    not extracted_page_text
                    or extracted_metrics["word_count"] < 12
                    or extracted_metrics["quality_score"] < 0.22
                )

                if not needs_ocr:
                    text_parts.append(f"Page {page_number}\n{extracted_page_text}")
                    continue

                pix = page.get_pixmap(matrix=fitz.Matrix(render_scale, render_scale), alpha=False)
                image = Image.open(BytesIO(pix.tobytes("png")))
                image = ImageOps.grayscale(image)
                image = ImageOps.autocontrast(image)
                image = image.filter(ImageFilter.SHARPEN)
                if pytesseract:
                    page_text = pytesseract.image_to_string(image, config="--psm 6")
                else:
                    import numpy as np
                    ocr_result, _ = rapid_ocr(np.array(image))
                    page_text = "\n".join(item[1] for item in ocr_result) if ocr_result else ""
                cleaned = cls.clean_ocr_text(page_text)
                if cleaned:
                    text_parts.append(f"Page {page_number}\n{cleaned}")

            full_text = cls.clean_text("\n\n".join(text_parts))
            chunk_size, overlap = cls._chunk_settings(
                file_size_bytes=len(file_content),
                page_count=pdf.page_count,
                used_ocr=True,
                text_length=len(full_text),
            )
            chunks = cls.chunk_text(full_text, chunk_size=chunk_size, overlap=overlap)
            quality_metrics = cls._text_quality_metrics(full_text)

            if not chunks:
                return {
                    "success": False,
                    "error": "OCR ran, but no readable text could be recognized from this scanned PDF.",
                    "chunks": []
                }

            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "used_ocr": True,
                "ocr_engine": "tesseract" if ocr_with_tesseract else "rapidocr",
                "ocr_quality_score": quality_metrics.get("quality_score", 0.0),
                "extraction_quality": quality_metrics,
                "chunk_size": chunk_size,
            }
        except Exception as exc:
            logger.error("OCR processing failed: %s", exc)
            return {
                "success": False,
                "error": f"OCR failed while processing the scanned PDF: {exc}",
                "chunks": []
            }
    
    @classmethod
    def process_docx(cls, file_content: bytes) -> Dict:
        """Extract text from DOCX file"""
        try:
            from docx import Document
            
            doc = Document(BytesIO(file_content))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract text from tables
            for table_index, table in enumerate(doc.tables, start=1):
                table_rows = []
                for row in table.rows:
                    table_rows.append([cell.text.strip() for cell in row.cells])
                text_parts.extend(
                    cls._format_table_matrix(table_rows, table_label=f"DOCX table {table_index}")
                )
            
            full_text = cls.clean_text("\n".join(text_parts))
            chunks = cls.chunk_text(full_text)
            
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "images": cls._extract_docx_images(file_content),
                "paragraph_count": len(doc.paragraphs)
            }
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            return {"success": False, "error": str(e), "chunks": []}
    
    @classmethod
    def process_txt(cls, file_content: bytes) -> Dict:
        """Extract text from TXT file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    text = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                text = file_content.decode('utf-8', errors='ignore')
            
            full_text = cls.clean_text(text)
            chunks = cls.chunk_text(full_text)
            
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks
            }
        except Exception as e:
            logger.error(f"Error processing TXT: {e}")
            return {"success": False, "error": str(e), "chunks": []}
    
    @classmethod
    def process_csv(cls, file_content: bytes) -> Dict:
        """Extract text from CSV file"""
        try:
            import csv
            from io import StringIO
            
            # Decode content
            text = file_content.decode('utf-8', errors='ignore')
            reader = csv.reader(StringIO(text))
            
            rows = []
            headers = None
            for i, row in enumerate(reader):
                if i == 0:
                    headers = row
                    rows.append(" | ".join(row))
                else:
                    # Create key-value pairs for better context
                    if headers:
                        row_text = ", ".join(f"{h}: {v}" for h, v in zip(headers, row) if v.strip())
                    else:
                        row_text = " | ".join(row)
                    if row_text:
                        rows.append(row_text)
            
            full_text = cls.clean_text("\n".join(rows))
            chunks = cls.chunk_text(full_text)
            
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "row_count": len(rows)
            }
        except Exception as e:
            logger.error(f"Error processing CSV: {e}")
            return {"success": False, "error": str(e), "chunks": []}

    @staticmethod
    def _json_get_case_insensitive(record: Dict, *keys):
        """Read common dataset keys even when Kaggle exports use User/AI or camelCase."""
        if not isinstance(record, dict):
            return None

        exact = next((record.get(key) for key in keys if key in record and record.get(key) is not None), None)
        if exact is not None:
            return exact

        lowered = {str(key).lower(): value for key, value in record.items()}
        return next((lowered.get(str(key).lower()) for key in keys if lowered.get(str(key).lower()) is not None), None)

    @classmethod
    def _flatten_json_value_to_lines(cls, value, label: str = "Text", lines=None, max_lines: int = 260):
        """Last-resort recursive extractor for readable JSON datasets with unknown schemas."""
        if lines is None:
            lines = []
        if len(lines) >= max_lines:
            return lines

        if isinstance(value, dict):
            for key, child in value.items():
                if len(lines) >= max_lines:
                    break
                cls._flatten_json_value_to_lines(child, str(key).replace("_", " ").title(), lines, max_lines)
            return lines

        if isinstance(value, list):
            for child in value:
                if len(lines) >= max_lines:
                    break
                cls._flatten_json_value_to_lines(child, label, lines, max_lines)
            return lines

        text_value = cls.clean_text(str(value or ""))
        if text_value and len(text_value) >= 2:
            lines.append(f"{label}: {text_value}")
        return lines

    @classmethod
    def _conversation_record_to_text(cls, record, index: int = 0) -> str:
        """Normalize common chatbot/conversation JSON shapes into readable text."""
        lines = []

        def add_line(role: str, value) -> None:
            text_value = cls.clean_text(str(value or ""))
            if text_value:
                lines.append(f"{role}: {text_value}")

        if isinstance(record, dict):
            title = cls._json_get_case_insensitive(record, "title", "topic", "category", "intent", "tag", "name")
            if title:
                lines.append(f"Conversation {index + 1} topic: {cls.clean_text(str(title))}")
            topic_line_count = len(lines)

            messages = cls._json_get_case_insensitive(
                record,
                "messages",
                "conversation",
                "conversations",
                "dialog",
                "dialogue",
                "turns",
                "chat",
                "transcript",
            )
            if isinstance(messages, list):
                for turn in messages:
                    if isinstance(turn, dict):
                        role = (
                            cls._json_get_case_insensitive(turn, "role", "from", "speaker", "sender", "author")
                            or "message"
                        )
                        value = cls._json_get_case_insensitive(
                            turn,
                            "content",
                            "text",
                            "value",
                            "message",
                            "utterance",
                            "human",
                            "user",
                            "assistant",
                            "bot",
                            "gpt",
                        )
                        if value:
                            add_line(str(role).title(), value)
                        else:
                            # Turns such as {"User": "...", "AI": "..."} carry role names as keys.
                            for key, child in turn.items():
                                if isinstance(child, (dict, list)):
                                    continue
                                add_line(str(key).replace("_", " ").title(), child)
                    else:
                        add_line("Message", turn)
            else:
                # Instruction/chat datasets usually use one of these pairings.
                pair_fields = [
                    ("Instruction", cls._json_get_case_insensitive(record, "instruction", "system")),
                    ("Input", cls._json_get_case_insensitive(record, "input", "context")),
                    ("User", cls._json_get_case_insensitive(record, "prompt", "question", "query", "user", "human")),
                    (
                        "Assistant",
                        cls._json_get_case_insensitive(
                            record,
                            "response",
                            "answer",
                            "assistant",
                            "output",
                            "completion",
                            "bot",
                            "gpt",
                            "ai",
                        ),
                    ),
                ]
                for role, value in pair_fields:
                    add_line(role, value)

            if len(lines) <= topic_line_count:
                # Last-resort recursive flattening keeps unknown JSON schemas searchable.
                lines.extend(cls._flatten_json_value_to_lines(record))

        elif isinstance(record, list):
            for item in record:
                lines.append(cls._conversation_record_to_text(item, index=index))
        else:
            add_line("Text", record)

        return cls.clean_text("\n".join(line for line in lines if line))

    @classmethod
    def _json_payload_to_text(cls, payload) -> str:
        """Extract conversation-oriented text from arbitrary JSON payloads."""
        records = payload
        if isinstance(payload, dict):
            for key in (
                "data",
                "intents",
                "conversations",
                "messages",
                "records",
                "items",
                "examples",
                "rows",
                "train",
                "validation",
                "paragraphs",
                "qas",
            ):
                if isinstance(payload.get(key), list):
                    records = payload[key]
                    break

        if isinstance(records, list):
            documents = []
            for index, record in enumerate(records):
                text = cls._conversation_record_to_text(record, index=index)
                if len(text) >= 40:
                    documents.append(text)
            return cls.clean_text("\n\n---\n\n".join(documents))

        return cls._conversation_record_to_text(records)

    @classmethod
    def process_json(cls, file_content: bytes) -> Dict:
        """Index JSON/JSONL datasets such as multi-turn chatbot conversations."""
        try:
            text = file_content.decode("utf-8", errors="ignore").strip()
            if not text:
                return {"success": False, "error": "JSON file is empty", "chunks": []}

            try:
                payload = json.loads(text)
                full_text = cls._json_payload_to_text(payload)
            except json.JSONDecodeError:
                # JSONL fallback.
                rows = []
                for index, line in enumerate(text.splitlines()):
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        rows.append(cls._conversation_record_to_text(json.loads(line), index=index))
                    except json.JSONDecodeError:
                        continue
                full_text = cls.clean_text("\n\n---\n\n".join(row for row in rows if row))

            chunk_size, overlap = cls._chunk_settings(
                file_size_bytes=len(file_content),
                text_length=len(full_text),
            )
            chunks = cls.chunk_text(full_text, chunk_size=chunk_size, overlap=overlap)
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "record_type": "conversation_json",
                "chunk_size": chunk_size,
            }
        except Exception as e:
            logger.error("Error processing JSON dataset: %s", e)
            return {"success": False, "error": str(e), "chunks": []}

    @classmethod
    def process_zip_json(cls, file_content: bytes) -> Dict:
        """Find JSON/JSONL files inside a ZIP export and index them together."""
        try:
            texts = []
            with zipfile.ZipFile(BytesIO(file_content)) as archive:
                names = [
                    name for name in archive.namelist()
                    if name.lower().endswith((".json", ".jsonl")) and not name.endswith("/")
                ]
                if not names:
                    return {"success": False, "error": "ZIP did not contain any JSON or JSONL files", "chunks": []}

                for name in names[:12]:
                    result = cls.process_json(archive.read(name))
                    if result.get("success") and result.get("text"):
                        texts.append(f"Dataset file: {name}\n{result['text']}")

            full_text = cls.clean_text("\n\n---\n\n".join(texts))
            chunk_size, overlap = cls._chunk_settings(
                file_size_bytes=len(file_content),
                text_length=len(full_text),
            )
            chunks = cls.chunk_text(full_text, chunk_size=chunk_size, overlap=overlap)
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "record_type": "conversation_zip_json",
                "chunk_size": chunk_size,
            }
        except Exception as e:
            logger.error("Error processing ZIP JSON dataset: %s", e)
            return {"success": False, "error": str(e), "chunks": []}
    
    @classmethod
    def process_pptx(cls, file_content: bytes) -> Dict:
        """Extract text from PPTX file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(BytesIO(file_content))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"Slide {slide_num}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                    elif hasattr(shape, "table"):
                        table_rows = []
                        for row in shape.table.rows:
                            table_rows.append([cell.text.strip() for cell in row.cells])
                        slide_texts.extend(
                            cls._format_table_matrix(table_rows, table_label=f"PPTX table on slide {slide_num}")
                        )
                
                if len(slide_texts) > 1:
                    text_parts.append("\n".join(slide_texts))
            
            full_text = cls.clean_text("\n\n".join(text_parts))
            chunks = cls.chunk_text(full_text)
            
            return {
                "success": True,
                "text": full_text,
                "chunks": chunks,
                "slide_count": len(prs.slides)
            }
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return {"success": False, "error": str(e), "chunks": []}
    
    @classmethod
    async def process_url(cls, url: str) -> Dict:
        """Scrape and extract text from a URL"""
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.get(url, follow_redirects=True)
                response.raise_for_status()
                content_type = response.headers.get("Content-Type", "").split(";", 1)[0].strip().lower()
                normalized_url = str(response.url).lower()

                if (
                    content_type in {"application/json", "text/json", "application/x-ndjson"}
                    or normalized_url.endswith((".json", ".jsonl"))
                ):
                    return cls.process_json(response.content)

                if content_type in {"application/zip", "application/x-zip-compressed"} or normalized_url.endswith(".zip"):
                    return cls.process_zip_json(response.content)
                
                soup = BeautifulSoup(response.text, 'lxml')
                title = soup.find('title')
                title_text = title.get_text(strip=True) if title else url
                images = cls.extract_html_images(soup, url, source_title=title_text)
                
                # Remove script and style elements
                for script in soup(["script", "style", "nav", "footer", "header"]):
                    script.decompose()
                
                # Get text from main content areas
                main_content = soup.find(['main', 'article']) or soup.find('body')
                
                if main_content:
                    # Extract paragraphs and headings
                    text_parts = []
                    for element in main_content.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'td']):
                        text = element.get_text(strip=True)
                        if text and len(text) > 10:
                            text_parts.append(text)
                    
                    full_text = cls.clean_text("\n".join(text_parts))
                else:
                    full_text = cls.clean_text(soup.get_text())
                
                chunks = cls.chunk_text(full_text)
                
                return {
                    "success": True,
                    "text": full_text,
                    "chunks": chunks,
                    "images": images,
                    "title": title_text,
                    "url": url
                }
        except Exception as e:
            logger.error(f"Error processing URL {url}: {e}")
            return {"success": False, "error": str(e), "chunks": []}
    
    @classmethod
    def process_file(cls, file_content: bytes, file_type: str) -> Dict:
        """Process file based on type"""
        normalized_file_type = file_type.lower()
        if normalized_file_type == "json" and zipfile.is_zipfile(BytesIO(file_content)):
            normalized_file_type = "zip"

        processors = {
            "pdf": cls.process_pdf,
            "docx": cls.process_docx,
            "txt": cls.process_txt,
            "csv": cls.process_csv,
            "json": cls.process_json,
            "zip": cls.process_zip_json,
            "pptx": cls.process_pptx,
        }
        
        processor = processors.get(normalized_file_type)
        if not processor:
            return {"success": False, "error": f"Unsupported file type: {file_type}", "chunks": []}
        
        result = processor(file_content)
        if result.get("success") and not result.get("chunks"):
            return {
                "success": False,
                "error": "No readable text could be extracted from this file.",
                "chunks": []
            }

        return result
